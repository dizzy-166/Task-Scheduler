# -*- coding: utf-8 -*-
"""Презентация к диплому «Поток» — приведение к методичке ГБПОУ НГК (15 слайдов, Arial, текст >=18 пт)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

TITLE_BLUE = RGBColor(0x2F, 0x5C, 0x9E)
ACCENT     = RGBColor(0x3B, 0x6F, 0xB0)
ACCENT_DK  = RGBColor(0x2B, 0x55, 0x8C)
DARK       = RGBColor(0x22, 0x22, 0x22)
GRAY       = RGBColor(0x55, 0x55, 0x55)
LIGHT      = RGBColor(0xEF, 0xF3, 0xFA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BD    = RGBColor(0xCF, 0xDC, 0xEE)
FONT = "Arial"

MEDIA = r"C:\Users\denis\OneDrive\Desktop\GithubClone\Task-Scheduler\Documentation\diploma_unpacked\word\media"
def img(n): return os.path.join(MEDIA, n)
IMG = {'login':img('image5.png'),'dashboard':img('image6.png'),'kanban':img('image7.png'),
    'task_modal':img('image8.png'),'gantt':img('image9.png'),'analytics':img('image10.png'),
    'chat':img('image11.png'),'roles':img('image12.png'),'ai_gen':img('image13.png'),
    'usecase':img('image14.png'),'state':img('image15.png'),'architecture':img('image16.png'),
    'seq_auth':img('image17.png'),'seq_status':img('image18.png'),'seq_create':img('image19.png'),
    'seq_chat':img('image20.png'),'deployment':img('image21.png'),'dataflow':img('image22.png'),
    'erd':img('image2.png'),'swagger':img('image23.png')}

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5; BLANK = prs.slide_layouts[6]; TOTAL = 15

def _rad(sp, f):
    try: sp.adjustments[0] = f
    except Exception: pass
def rect(slide,x,y,w,h,fill=None,line=None,line_w=None,shape=MSO_SHAPE.RECTANGLE,radius=None,shadow=False):
    sp=slide.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    if radius is not None: _rad(sp,radius)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(line_w or 1)
    sp.shadow.inherit=False
    if shadow:
        el=sp._element.spPr; ef=el.makeelement(qn('a:effectLst'),{}); el.append(ef)
        sh=ef.makeelement(qn('a:outerShdw'),{'blurRad':'90000','dist':'40000','dir':'5400000','rotWithShape':'0'}); ef.append(sh)
        c=sh.makeelement(qn('a:srgbClr'),{'val':'1A3A66'}); sh.append(c)
        a=c.makeelement(qn('a:alpha'),{'val':'20000'}); c.append(a)
    return sp
def _runs(p,parts,size,color,bold,font):
    if isinstance(parts,str): parts=[(parts,{})]
    for t,ov in parts:
        r=p.add_run(); r.text=t; f=r.font
        f.size=Pt(ov.get('size',size)); f.name=ov.get('font',font)
        f.bold=ov.get('bold',bold); f.italic=ov.get('italic',False); f.color.rgb=ov.get('color',color)
def text(slide,x,y,w,h,lines,size=18,color=DARK,bold=False,font=FONT,align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP,line_spacing=1.0,sp_after=3,wrap=True):
    tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
    if isinstance(lines,(str,tuple)): lines=[lines]
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if line_spacing: p.line_spacing=line_spacing
        p.space_after=Pt(sp_after); p.space_before=Pt(0)
        _runs(p,ln,size,color,bold,font)
    return tb
def bullets(slide,x,y,w,h,items,size=18,gap=7,line_spacing=1.05):
    tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); lvl=it.get('lvl',0)
        p.line_spacing=line_spacing; p.space_after=Pt(it.get('gap',gap)); p.space_before=Pt(0)
        sz=it.get('size',size-(2 if lvl else 0)); col=it.get('c',DARK if lvl else TITLE_BLUE); bold=it.get('b',not lvl)
        marker=it.get('marker', '' if it.get('nomark') else ('—  ' if lvl else '•  '))
        indent=it.get('indent',0.32*lvl)
        pPr=p._p.get_or_add_pPr(); pPr.set('marL',str(int(Inches(indent+(0.3 if marker else 0))))); pPr.set('indent',str(int(Inches(-0.3 if marker else 0))))
        parts=it.get('parts')
        if parts is None: parts=[(marker+it['t'],{})] if marker else [(it['t'],{})]
        elif marker: parts=[(marker,{})]+parts
        for t,ov in parts:
            r=p.add_run(); r.text=t; f=r.font
            f.size=Pt(ov.get('size',sz)); f.name=ov.get('font',FONT)
            f.bold=ov.get('bold',bold); f.italic=ov.get('italic',False); f.color.rgb=ov.get('color',col)
    return tb
def pic_fit(slide,path,bx,by,bw,bh,align='center',valign='middle',frame=True,shadow=True):
    iw,ih=Image.open(path).size; ar=iw/ih; bar=bw/bh
    if ar>bar: w=bw; h=bw/ar
    else: h=bh; w=bh*ar
    x=bx+(bw-w)/2 if align=='center' else (bx if align=='left' else bx+(bw-w))
    y=by+(bh-h)/2 if valign=='middle' else (by if valign=='top' else by+(bh-h))
    if frame: rect(slide,x-0.04,y-0.04,w+0.08,h+0.08,fill=WHITE,line=CARD_BD,line_w=1,radius=0.03,shadow=shadow)
    return slide.shapes.add_picture(path,Inches(x),Inches(y),Inches(w),Inches(h))
def brand(slide):
    rect(slide,0.42,0.34,0.46,0.46,fill=ACCENT,radius=0.28)
    rect(slide,0.53,0.45,0.24,0.045,fill=WHITE,radius=0.5)
    rect(slide,0.53,0.545,0.17,0.045,fill=WHITE,radius=0.5)
    rect(slide,0.53,0.64,0.21,0.045,fill=WHITE,radius=0.5)
    text(slide,0.96,0.33,2.0,0.5,"Поток",size=20,bold=True,color=TITLE_BLUE,anchor=MSO_ANCHOR.MIDDLE)
def page(slide,n):
    tb=slide.shapes.add_textbox(Inches(11.05),Inches(6.92),Inches(2.0),Inches(0.45)); tf=tb.text_frame; tf.word_wrap=False
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
    _runs(p,[(f"{n:02d}",{'bold':True,'size':18,'color':ACCENT}),(" / ",{'size':14,'color':GRAY}),(f"{TOTAL}",{'size':13,'color':GRAY})],18,ACCENT,True,FONT)
def header(slide,n,title,title_size=30,tx=2.0,tw=9.33):
    brand(slide)
    text(slide,tx,0.34,tw,0.85,title,size=title_size,bold=True,color=TITLE_BLUE,anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
    cx=tx+tw/2-0.525; rect(slide,cx,1.16,1.05,0.07,fill=ACCENT,radius=0.5); page(slide,n)
def content(n,title,**kw):
    s=prs.slides.add_slide(BLANK); rect(s,0,0,SW,SH,fill=WHITE); header(s,n,title,**kw); return s
def caption(slide,x,y,w,txt):
    text(slide,x,y,w,0.32,txt,size=15,color=GRAY,align=PP_ALIGN.CENTER)
def table(slide,x,y,w,rows,col_w,header_h=0.5,row_h=0.5,fs=16,hfs=16):
    nrow=len(rows); ncol=len(rows[0])
    gt=slide.shapes.add_table(nrow,ncol,Inches(x),Inches(y),Inches(w),Inches(header_h+row_h*(nrow-1))).table
    gt._tbl.tblPr.set('firstRow','0'); gt._tbl.tblPr.set('bandRow','0')
    for j,cw in enumerate(col_w): gt.columns[j].width=Inches(cw)
    gt.rows[0].height=Inches(header_h)
    for i in range(1,nrow): gt.rows[i].height=Inches(row_h)
    for i,row in enumerate(rows):
        for j,val in enumerate(row):
            c=gt.cell(i,j); c.margin_left=Inches(0.08); c.margin_right=Inches(0.06)
            c.margin_top=Inches(0.01); c.margin_bottom=Inches(0.01); c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.fill.solid(); c.fill.fore_color.rgb=ACCENT_DK if i==0 else (WHITE if i%2 else LIGHT)
            tf=c.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]
            p.alignment=PP_ALIGN.LEFT if (j==0 and i>0) else PP_ALIGN.CENTER
            parts=val if isinstance(val,list) else [(str(val),{})]
            for txt,ov in parts:
                r=p.add_run(); r.text=txt; f=r.font; f.name=FONT
                f.size=Pt(ov.get('size',hfs if i==0 else fs)); f.bold=ov.get('bold',i==0 or j==0)
                f.color.rgb=ov.get('color',WHITE if i==0 else DARK)
    return gt
YES=[("✓",{'color':RGBColor(0x1E,0x9E,0x57),'bold':True,'size':19})]
NO =[("—",{'color':RGBColor(0xC0,0x3A,0x3A),'bold':True,'size':19})]
PART=[("~",{'color':RGBColor(0xD9,0x8A,0x10),'bold':True,'size':19})]
print("helpers ready (Arial)")

# 1 — Титульный
s=prs.slides.add_slide(BLANK); rect(s,0,0,SW,SH,fill=WHITE)
rect(s,0,0,SW,2.1,fill=TITLE_BLUE); rect(s,0,2.1,SW,0.10,fill=ACCENT)
text(s,0.8,0.42,11.7,1.4,["Министерство образования Нижегородской области",
 "Государственное бюджетное профессиональное образовательное учреждение",
 "«Нижегородский Губернский колледж»"],size=15,color=WHITE,align=PP_ALIGN.CENTER,line_spacing=1.12,sp_after=2)
rect(s,5.87,2.55,0.95,0.95,fill=ACCENT,radius=0.26,shadow=True)
rect(s,6.09,2.78,0.5,0.095,fill=WHITE,radius=0.5); rect(s,6.09,2.97,0.35,0.095,fill=WHITE,radius=0.5); rect(s,6.09,3.16,0.44,0.095,fill=WHITE,radius=0.5)
text(s,1.0,3.72,11.33,0.5,"ДИПЛОМНАЯ РАБОТА",size=18,bold=True,color=ACCENT_DK,align=PP_ALIGN.CENTER)
text(s,1.2,4.2,10.93,1.0,"Проектирование и разработка планировщика задач для компании",
 size=27,bold=True,color=TITLE_BLUE,align=PP_ALIGN.CENTER,line_spacing=1.04)
text(s,1.0,5.28,11.33,0.4,"Корпоративная система управления задачами «Поток»",size=16,color=GRAY,align=PP_ALIGN.CENTER)
rect(s,5.17,5.78,3.0,0.04,fill=CARD_BD)
text(s,1.0,5.86,11.33,0.35,"Специальность 09.02.07 Информационные системы и программирование",size=14,color=GRAY,align=PP_ALIGN.CENTER)
rect(s,7.18,6.36,5.45,0.95,fill=LIGHT,line=CARD_BD,line_w=1,radius=0.06)
bx=s.shapes.add_textbox(Inches(7.4),Inches(6.42),Inches(5.1),Inches(0.85)); tf=bx.text_frame; tf.word_wrap=True
for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
for i,r in enumerate([[("Выполнил: ",{'bold':True,'color':TITLE_BLUE}),("Шеронов Д. С., группа 41П",{})],
                      [("Руководитель: ",{'bold':True,'color':TITLE_BLUE}),("______________",{})]]):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(3); p.line_spacing=1.1
    for t,ov in r:
        rr=p.add_run(); rr.text=t; f=rr.font; f.name=FONT; f.size=Pt(18); f.bold=ov.get('bold',False); f.color.rgb=ov.get('color',DARK)
text(s,0.73,6.55,3.6,0.6,["Нижний Новгород","2026 г."],size=15,bold=True,color=TITLE_BLUE,line_spacing=1.05,sp_after=0)

# 2 — Актуальность
s=content(2,"Актуальность")
cards=[("Цифровизация труда","Качество инструментов управления задачами напрямую определяет успех проектной работы."),
 ("Удалённая работа","Распределённые команды сделали цифровую координацию критической инфраструктурой."),
 ("Импортозамещение","Указ Президента РФ № 166 от 30.03.2022 — переход на отечественное ПО вместо Jira и Asana."),
 ("Разрыв на рынке","Яндекс Трекер, Кайтен, YouGile не закрывают ИИ-интеграцию, RBAC и встроенный мессенджер.")]
cw,ch,gp=5.75,2.35,0.42; x0,y0=1.18,1.6
for i,(h,b) in enumerate(cards):
    cx=x0+(i%2)*(cw+gp); cy=y0+(i//2)*(ch+0.35)
    rect(s,cx,cy,cw,ch,fill=LIGHT,line=CARD_BD,line_w=1,radius=0.05,shadow=True); rect(s,cx,cy,0.13,ch,fill=ACCENT)
    text(s,cx+0.32,cy+0.22,cw-0.55,0.5,h,size=20,bold=True,color=TITLE_BLUE)
    text(s,cx+0.32,cy+0.82,cw-0.6,ch-1.0,b,size=16,color=DARK,line_spacing=1.12)

# 3 — Объект, предмет, цель
s=content(3,"Объект, предмет и цель исследования",title_size=28)
blk=[("Объект исследования","Корпоративные системы управления задачами как класс прикладного программного обеспечения.",1.18),
 ("Предмет исследования","Методы проектирования и разработки веб-системы управления задачами с ролевой моделью, канбан-визуализацией и интеграцией ИИ.",1.55),
 ("Цель работы","Проектирование и разработка корпоративной системы управления задачами «Поток» с поддержкой ролевой модели, визуализации данных и интеграции с ИИ-сервисами.",1.7)]
y=1.62
for h,b,bh in blk:
    rect(s,1.18,y,10.95,bh,fill=LIGHT,line=CARD_BD,line_w=1,radius=0.05,shadow=True); rect(s,1.18,y,0.14,bh,fill=ACCENT)
    text(s,1.5,y+0.16,10.4,0.4,h,size=20,bold=True,color=TITLE_BLUE)
    text(s,1.5,y+0.66,10.35,bh-0.78,b,size=17,color=DARK,line_spacing=1.12)
    y+=bh+0.3

# 4 — Задачи
s=content(4,"Задачи работы")
tasks=["Проанализировать предметную область и выявить функциональный разрыв",
 "Спроектировать архитектуру веб-приложения «клиент–сервер»",
 "Разработать гибкую ролевую модель доступа (RBAC)",
 "Реализовать канбан-доску с перетаскиванием карточек",
 "Интегрировать ИИ для генерации задач и аналитики",
 "Обеспечить экспорт отчётности в форматы PDF и Excel",
 "Реализовать встроенный мессенджер",
 "Развернуть приложение в облаке и провести тестирование"]
x0,y0,cw,rh=1.18,1.55,10.95,0.66
for i,t in enumerate(tasks):
    cy=y0+i*(rh+0.04); rect(s,x0,cy,cw,rh,fill=(WHITE if i%2 else LIGHT),line=CARD_BD,line_w=0.75,radius=0.06)
    rect(s,x0+0.12,cy+0.12,0.42,0.42,fill=ACCENT,radius=0.22)
    text(s,x0+0.12,cy+0.12,0.42,0.42,f"{i+1}",size=18,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,x0+0.74,cy,cw-0.9,rh,t,size=18,color=DARK,anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.0)

# 5 — Класс ПО + методология
s=content(5,"Планировщик задач как класс ПО")
rect(s,1.18,1.55,5.4,5.15,fill=LIGHT,line=CARD_BD,line_w=1,radius=0.04,shadow=True)
text(s,1.5,1.72,4.8,0.4,"Эволюция инструментов",size=19,bold=True,color=TITLE_BLUE)
evo=[("До 1980-х","Бумажные доски, канбан Toyota"),("1980–2000-е","Excel, MS Project"),
 ("2000–2015","SaaS: Jira, Trello, Basecamp"),("2015 — н.в.","ИИ-платформы и мессенджеры")]
yy=2.34
for tag,d in evo:
    rect(s,1.5,yy+0.04,1.6,0.46,fill=ACCENT,radius=0.2)
    text(s,1.5,yy+0.04,1.6,0.46,tag,size=14,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,3.25,yy,3.15,0.85,d,size=16,color=DARK,anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.0); yy+=1.05
rect(s,6.73,1.55,5.4,5.15,fill=LIGHT,line=CARD_BD,line_w=1,radius=0.04,shadow=True)
text(s,7.05,1.72,4.8,0.4,"Выбор методологии",size=19,bold=True,color=TITLE_BLUE)
for i,(m,d) in enumerate([("Waterfall","низкая адаптивность"),("Agile","нужна зрелая команда"),("Scrum","избыточен для малых команд")]):
    text(s,7.05,2.3+i*0.5,4.8,0.5,[[("•  "+m+"  ",{'bold':True,'color':DARK,'size':17}),("— "+d,{'color':GRAY,'size':15})]],size=17)
rect(s,7.05,4.0,4.78,2.55,fill=WHITE,line=ACCENT,line_w=1.5,radius=0.05); rect(s,7.05,4.0,4.78,0.6,fill=ACCENT,radius=0.05)
text(s,7.05,4.0,4.78,0.6,"Kanban — основа системы",size=18,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
bullets(s,7.28,4.78,4.4,1.7,[{'t':'низкий порог входа','lvl':1,'size':16,'gap':4},
 {'t':'наглядная визуализация','lvl':1,'size':16,'gap':4},
 {'t':'гибкость без спринтов','lvl':1,'size':16,'gap':4},
 {'t':'для разработки и операций','lvl':1,'size':16,'gap':4}],line_spacing=1.04)

# 6 — Анализ аналогов
s=content(6,"Сравнительный анализ аналогов",title_size=28)
rows=[["Система","Канбан","RBAC","Мессен-\nджер","ИИ-\nзадачи","Мульти-\nкомпания"],
 ["Jira",YES,YES,NO,PART,NO],["Trello",YES,NO,NO,NO,NO],["Asana",YES,NO,NO,PART,NO],
 ["Bitrix24",YES,PART,YES,NO,NO],["Яндекс Трекер",YES,NO,NO,NO,NO],["Кайтен",YES,NO,NO,NO,NO],
 ["YouGile",YES,NO,NO,NO,NO],[[("«Поток»",{'color':WHITE})],YES,YES,YES,YES,YES]]
gt=table(s,1.18,1.6,10.95,rows,[2.75,1.62,1.62,1.66,1.62,1.68],header_h=0.66,row_h=0.45,fs=16,hfs=14)
for j in range(6):
    c=gt.cell(8,j); c.fill.solid(); c.fill.fore_color.rgb=ACCENT
    if j==0:
        for r in c.text_frame.paragraphs[0].runs: r.font.color.rgb=WHITE; r.font.bold=True
rect(s,1.18,6.42,10.95,0.62,fill=LIGHT,line=CARD_BD,line_w=1,radius=0.06)
text(s,1.45,6.42,10.45,0.62,[[("Вывод:  ",{'bold':True,'color':TITLE_BLUE,'size':16}),
 ("ни одно решение не покрывает профиль полностью — дефицит RBAC, мессенджера и ИИ-генерации задач.",{'color':DARK,'size':16})]],size=16,anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.02)

# 7 — Требования
s=content(7,"Требования к системе")
text(s,1.18,1.5,5.4,0.4,"Роли пользователей",size=19,bold=True,color=TITLE_BLUE)
roles=[("Руководитель","аналитика, управление ролями, отчёты"),("Проектный менеджер","распределение задач, контроль сроков"),
 ("Исполнитель","свои задачи, комментарии, таймер"),("Внешний участник","ограниченный доступ через RBAC")]
yy=2.0
for r,d in roles:
    rect(s,1.18,yy,5.4,1.06,fill=LIGHT,line=CARD_BD,line_w=1,radius=0.05); rect(s,1.18,yy,0.12,1.06,fill=ACCENT)
    text(s,1.45,yy+0.13,5.0,0.4,r,size=18,bold=True,color=TITLE_BLUE)
    text(s,1.45,yy+0.56,5.0,0.45,d,size=15,color=DARK,line_spacing=1.02); yy+=1.17
rect(s,6.78,1.5,5.35,5.28,fill=WHITE,line=ACCENT,line_w=1.5,radius=0.04,shadow=True); rect(s,6.78,1.5,5.35,0.62,fill=ACCENT,radius=0.04)
text(s,6.78,1.5,5.35,0.62,"Эталонный профиль",size=19,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
bullets(s,7.05,2.36,4.85,3.5,[{'t':'Покрытие бизнес-процессов','lvl':1,'size':17},
 {'t':'Гибкая ролевая модель (RBAC)','lvl':1,'size':17},{'t':'Встроенный мессенджер','lvl':1,'size':17},
 {'t':'Интеграция с ИИ','lvl':1,'size':17},{'t':'Производительность и безопасность','lvl':1,'size':17},
 {'t':'Ориентация на рынок РФ','lvl':1,'size':17}],gap=12,line_spacing=1.08)
rect(s,7.05,5.98,4.8,0.62,fill=LIGHT,line=CARD_BD,line_w=1,radius=0.06)
text(s,7.2,5.98,4.55,0.62,[[("Стандарты:  ",{'bold':True,'color':TITLE_BLUE,'size':15}),
 ("ГОСТ 34.602-2020, 25010-2015",{'color':DARK,'size':15})]],size=15,anchor=MSO_ANCHOR.MIDDLE)

# 8 — Технологический стек
s=content(8,"Технологический стек",title_size=30)
techs=[("Бэкенд","Django + DRF"),("Фронтенд","React 19 + Vite"),("Состояние","Zustand"),
 ("СУБД","PostgreSQL 16"),("ИИ","Cerebras · Llama 3.1"),("Деплой","Render + Vercel"),
 ("Безопасность","JWT · PBKDF2 · HTTPS"),("Контроль версий","Git · GitHub")]
cw,ch,gx,gy=5.42,1.0,0.3,0.18; x0,y0=1.18,1.55
for i,(cat,tech) in enumerate(techs):
    cx=x0+(i%2)*(cw+gx); cy=y0+(i//2)*(ch+gy)
    rect(s,cx,cy,cw,ch,fill=LIGHT,line=CARD_BD,line_w=1,radius=0.06); rect(s,cx,cy,0.13,ch,fill=ACCENT)
    text(s,cx+0.32,cy+0.13,cw-0.5,0.4,cat,size=15,bold=True,color=GRAY)
    text(s,cx+0.32,cy+0.46,cw-0.5,0.45,tech,size=19,bold=True,color=TITLE_BLUE)
text(s,1.18,6.66,10.95,0.4,[[("Импортозамещение:  ",{'bold':True,'color':TITLE_BLUE,'size':15}),
 ("открытый стек и модель Llama допускают локальное развёртывание.",{'color':DARK,'size':15})]],size=15,anchor=MSO_ANCHOR.MIDDLE)

print("slides 1-8 ok")

# 9 — Проектирование: архитектура + БД (2 рисунка)
s=content(9,"Проектирование системы")
pic_fit(s,IMG['architecture'],0.55,1.55,6.2,4.55,valign='top')
caption(s,0.55,6.2,6.2,"Рисунок 1 — Архитектура «клиент–сервер»")
pic_fit(s,IMG['erd'],6.95,1.55,6.05,4.55,valign='top')
caption(s,6.95,6.2,6.05,"Рисунок 2 — ER-диаграмма базы данных (3НФ, UUID)")

# 10 — Аутентификация и ролевая модель
s=content(10,"Аутентификация и ролевая модель")
bullets(s,0.75,1.75,4.45,4.7,[{'t':'JWT-аутентификация','c':TITLE_BLUE},
 {'t':'пара токенов access + refresh','lvl':1},{'t':'автообновление без выхода','lvl':1},
 {'t':'Безопасность','c':TITLE_BLUE},{'t':'хэширование паролей PBKDF2','lvl':1},
 {'t':'верификация email','lvl':1},{'t':'Ролевая модель RBAC','c':TITLE_BLUE},
 {'t':'гранулированные разрешения','lvl':1},{'t':'класс CanManageTask','lvl':1}],size=18,gap=8,line_spacing=1.06)
pic_fit(s,IMG['seq_auth'],5.25,1.55,7.75,4.85)
caption(s,5.25,6.52,7.75,"Рисунок 3 — Диаграмма последовательности: вход и выдача JWT")

# 11 — Канбан-доска
s=content(11,"Канбан-доска и управление задачами",title_size=27)
pic_fit(s,IMG['kanban'],0.7,1.55,8.05,4.35,valign='top')
caption(s,0.7,6.0,8.05,"Рисунок 4 — Канбан-доска системы «Поток»")
rect(s,8.95,1.55,4.08,4.95,fill=LIGHT,line=CARD_BD,line_w=1,radius=0.04,shadow=True)
bullets(s,9.2,1.8,3.65,4.6,[{'t':'Drag-and-drop','c':TITLE_BLUE},{'t':'автосмена статуса','lvl':1,'size':16},
 {'t':'Настраиваемые колонки','c':TITLE_BLUE},{'t':'Карточка задачи','c':TITLE_BLUE},
 {'t':'подзадачи, таймер','lvl':1,'size':16},{'t':'Фильтрация и поиск','c':TITLE_BLUE},
 {'t':'Мягкое удаление','c':TITLE_BLUE}],size=18,gap=10,line_spacing=1.05)

# 12 — Генерация задач ИИ
s=content(12,"Генерация задач с помощью ИИ")
bullets(s,0.75,1.8,4.55,4.6,[{'t':'Cerebras · Llama 3.1','c':TITLE_BLUE},
 {'t':'модель 8B параметров','lvl':1},{'t':'быстрый инференс','lvl':1},
 {'t':'Как это работает','c':TITLE_BLUE},{'t':'ввод описания проекта','lvl':1},
 {'t':'ИИ предлагает список задач','lvl':1},{'t':'с приоритетами и сроками','lvl':1},
 {'t':'Эффект','c':TITLE_BLUE},{'t':'декомпозиция за минуты','lvl':1}],size=18,gap=8,line_spacing=1.06)
pic_fit(s,IMG['ai_gen'],5.7,1.55,7.0,4.95)
caption(s,5.7,6.56,7.0,"Рисунок 5 — Генератор задач с использованием ИИ")

# 13 — Аналитика и мессенджер (2 рисунка)
s=content(13,"Аналитика, отчётность и коммуникации",title_size=26)
pic_fit(s,IMG['analytics'],0.55,1.7,6.2,3.5,valign='top')
caption(s,0.55,5.3,6.2,"Рисунок 6 — Модуль аналитики (PDF / Excel)")
pic_fit(s,IMG['chat'],6.95,1.7,6.05,3.5,valign='top')
caption(s,6.95,5.3,6.05,"Рисунок 7 — Встроенный мессенджер")
text(s,1.0,5.85,11.3,0.7,[[("Аналитика: ",{'bold':True,'color':TITLE_BLUE,'size':16}),
 ("метрики, диаграммы, ИИ-рекомендации, экспорт.   ",{'color':DARK,'size':16}),
 ("Мессенджер: ",{'bold':True,'color':TITLE_BLUE,'size':16}),
 ("корпоративный, проектные и личные каналы, @mention.",{'color':DARK,'size':16})]],
 size=16,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.1)

# 14 — Тестирование и внедрение
s=content(14,"Тестирование и внедрение")
text(s,1.18,1.5,5.5,0.4,"Тестирование",size=19,bold=True,color=TITLE_BLUE)
bullets(s,1.2,2.0,5.3,2.7,[{'t':'unit, интеграционные, end-to-end','lvl':1,'size':17},
 {'t':'10 тест-сценариев пройдено','lvl':1,'size':17},{'t':'API — Postman, UI — 3 браузера','lvl':1,'size':17},
 {'t':'PEP 8, ORM, ESLint','lvl':1,'size':17},{'t':'дефекты устранены и перепроверены','lvl':1,'size':17}],gap=9,line_spacing=1.08)
rect(s,1.18,4.55,5.35,2.05,fill=LIGHT,line=CARD_BD,line_w=1,radius=0.05)
text(s,1.45,4.72,4.9,0.4,"Развёрнуто в облаке",size=18,bold=True,color=TITLE_BLUE)
bullets(s,1.45,5.22,4.95,1.3,[{'t':'бэкенд — Render, фронтенд — Vercel','lvl':1,'size':16},
 {'t':'CI/CD: автодеплой из GitHub','lvl':1,'size':16},
 {'t':'task-scheduler-snowy.vercel.app','lvl':1,'size':16,'c':ACCENT_DK,'b':True}],gap=6,line_spacing=1.05)
pic_fit(s,IMG['deployment'],6.75,1.6,6.3,4.9)
caption(s,6.75,6.55,6.3,"Рисунок 8 — Диаграмма развёртывания")

# 15 — Заключение
s=content(15,"Заключение")
res=[("Разработана система «Поток»","полнофункциональное веб-приложение корпоративного класса"),
 ("REST API (40+ эндпоинтов)","JWT-аутентификация и ролевая модель RBAC"),
 ("SPA-интерфейс","канбан, Ганта, мессенджер, аналитика, @mention"),
 ("Интеграция с ИИ","генерация задач и аналитика на базе Llama 3.1"),
 ("Развёрнута в облаке","Render + Vercel, CI/CD, прошла тестирование")]
yy=1.62
for h,d in res:
    rect(s,1.18,yy,11.0,0.85,fill=LIGHT,line=CARD_BD,line_w=1,radius=0.06); rect(s,1.18,yy,0.13,0.85,fill=ACCENT)
    rect(s,1.45,yy+0.22,0.4,0.4,fill=ACCENT,radius=0.5)
    text(s,1.45,yy+0.21,0.4,0.4,"✓",size=16,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    text(s,2.05,yy+0.11,10.0,0.4,h,size=18,bold=True,color=TITLE_BLUE)
    text(s,2.05,yy+0.5,10.0,0.32,d,size=15,color=DARK); yy+=0.95
rect(s,1.18,yy+0.02,11.0,0.62,fill=ACCENT,radius=0.06)
text(s,1.18,yy+0.02,11.0,0.62,"Цель достигнута — система готова к эксплуатации. Спасибо за внимание!",
 size=18,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

print("slides 9-15 ok")
OUT=r"C:\Users\denis\OneDrive\Desktop\Шеронов Денис Сергеевич.pptx"
prs.save(OUT)
print("SAVED", OUT, "| slides:", len(prs.slides._sldIdLst))
